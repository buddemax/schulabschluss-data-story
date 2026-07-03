# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Emu
p=Presentation(r"C:\Users\maxot\OneDrive\Desktop\datastory_school\Schulabschluss_DataStory_Praesentation.pptx")
SW,SH=p.slide_width,p.slide_height
for i,s in enumerate(p.slides,1):
    print(f"\n## Folie {i}")
    npic=0; oob=[]
    for sh in s.shapes:
        if sh.shape_type==13: npic+=1
        # out-of-bounds check
        try:
            if sh.left is not None and (sh.left<0 or sh.top<0 or sh.left+sh.width>SW+10000 or sh.top+sh.height>SH+10000):
                oob.append(getattr(sh,'name','?'))
        except: pass
        if sh.has_text_frame and sh.text_frame.text.strip():
            t=sh.text_frame.text.strip().replace("\n"," | ")
            print("   T:",t[:95])
    print(f"   [Bilder: {npic}]" + (f"  OOB:{oob}" if oob else ""))
